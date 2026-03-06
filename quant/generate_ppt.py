#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成量化交易系统PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
TITLE_COLOR = RGBColor(0x1a, 0x56, 0xdb)  # 蓝色
DARK_COLOR = RGBColor(0x33, 0x33, 0x33)   # 深灰色
ACCENT_COLOR = RGBColor(0xe6, 0x69, 0x00)  # 橙色

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets=None, two_columns=None):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    if two_columns:
        # 左侧内容
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        
        for i, item in enumerate(two_columns[0]):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(10)
        
        # 右侧内容
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        
        for i, item in enumerate(two_columns[1]):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(10)
    
    elif bullets:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, table_data, col_widths=None):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 表格
    rows = len(table_data)
    cols = len(table_data[0])
    
    if col_widths is None:
        col_widths = [Inches(12.333/cols)] * cols
    
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                     Inches(12.333), Inches(rows * 0.5)).table
    
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            cell_obj = table.cell(i, j)
            cell_obj.text = str(cell)
            
            if i == 0:  # 表头
                cell_obj.fill.solid()
                cell_obj.fill.fore_color.rgb = RGBColor(0x1a, 0x56, 0xdb)
                para = cell_obj.text_frame.paragraphs[0]
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                para.alignment = PP_ALIGN.CENTER
            else:
                para = cell_obj.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.font.color.rgb = DARK_COLOR
                para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    
    # 设置列宽
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    return slide

# ==================== 开始生成PPT ====================

# 1. 封面
add_title_slide(prs, "量化交易系统", "期货3连K策略 + A股关键点选股")

# 2. 目录
add_content_slide(prs, "目录", [
    "一、项目概述与系统架构",
    "二、期货3连K策略v3详解",
    "三、A股关键点选股策略",
    "四、回测实验结果",
    "五、实践生产与应用",
    "六、总结与展望"
])

# 3. 项目概述
add_content_slide(prs, "一、项目概述", [
    "构建完整的期货+A股量化交易系统",
    "涵盖数据获取、策略研发、回测验证到实盘交易全链路",
    "数据源: Tushare Pro API",
    "编程语言: Python 3.13",
    "消息通知: 飞书Webhook推送"
])

# 4. 系统架构
add_content_slide(prs, "系统架构", [
    "数据层: Tushare Pro + 本地CSV存储 + 实时数据流",
    "策略层: 期货3连K策略 + A股关键点选股策略",
    "回测层: 历史数据回测 + 胜率统计 + 收益分析",
    "应用层: 实时信号检测 + 交易提醒 + 监控面板"
])

# 5. 期货3连K策略 - 原理
add_content_slide(prs, "二、期货3连K策略v3 - 策略原理", [
    "核心理念: 放量K线突破 + 2根K线确认 + 第3根K线买入",
    "",
    "Step 1 - 放量K检测: 成交量 > 3倍20日均量",
    "Step 2 - 2K确认: 等待2根K线验证趋势延续",
    "Step 3 - 突破买入: 第3根K线突破前高时买入",
    "",
    "风控策略: 半仓1:1止盈止损，全仓2:1止盈止损"
])

# 6. 期货3连K策略 - 筛选
add_content_slide(prs, "二、期货3连K策略v3 - 主力合约筛选", [
    "筛选逻辑: 基于成交量最大的合约作为主力合约",
    "筛选范围: 6个月内到期合约",
    "",
    "交易所分布:",
    "• CFFEX (中金所): 21个 - IF/IC/IH/IM/T/TF/TS/TU",
    "• SHFE (上期所): 9个 - cu/al/zn/pb/ni/au/ag/rb/hc",
    "• DCE (大商所): 11个 - a/b/c/m/y/p/j/l/v/eg/pg",
    "• CZCE (郑商所): 9个 - c/cs/rm/ma/ta/zc/pf/jr/sr",
    "• INE (上能所): 3个 - sc/nr/bc",
    "• GFEX (广期所): 1个 - df",
    "",
    "总计: 54个主力合约，91个交易标的"
])

# 7. 期货回测结果
table_data = [
    ["指标", "数值"],
    ["总交易次数", "5,742笔"],
    ["盈利次数", "3,596笔"],
    ["亏损次数", "2,146笔"],
    ["胜率", "62.68%"],
    ["总收益", "+68,250元"],
    ["平均每笔收益", "+11.89元"],
]
add_table_slide(prs, "二、期货3连K策略v3 - 回测结果(2025.06-2025.12)", table_data)

# 8. 期货多空对比
table_data2 = [
    ["方向", "交易次数", "胜率", "收益"],
    ["做多", "3,896", "65.2%", "+48,550元"],
    ["做空", "1,846", "56.3%", "+19,700元"],
]
add_table_slide(prs, "二、期货3连K策略v3 - 多空对比", table_data2)

# 9. A股关键点选股 - 原理
add_content_slide(prs, "三、A股关键点选股 - 选股条件", [
    "核心理念: 在趋势启动点买入，等待价值回归",
    "",
    "【买入条件】(全部满足)",
    "• 近15日涨幅 < 15%",
    "• 当前价格距60日高点 < 20%",
    "• 成交量 > 1.5倍20日均量",
    "• 前2日中有1日成交量 < 20日均量(缩量确认)",
    "",
    "【剔除条件】",
    "• 创业板(300xxx) • 科创板(688xxx)",
    "• ST股票 • 北交所 • 市值<50亿元"
])

# 10. A股卖点分析
add_content_slide(prs, "三、A股关键点选股 - 卖点分析", [
    "【5%止盈止损策略】",
    "• 止盈: 买入后涨幅达到5%，卖出锁定利润",
    "• 止损: 买入后跌幅达到5%，卖出控制风险",
    "• 持有: 30天后强制平仓",
    "",
    "【卖点分析指标】",
    "• 21日最大收益/回撤",
    "• 30日最大收益/回撤",
    "• 最佳卖点日期",
    "• 止盈止损统计"
])

# 11. A股回测结果
table_data3 = [
    ["指标", "数值"],
    ["分析股票", "3,518只"],
    ["有买点的股票", "3,478只"],
    ["买点总数", "28,336个"],
    ["胜率(止盈/止盈+止损)", "59.6%"],
    ["平均21日收益", "+9.8%"],
    ["中位数21日收益", "+5.7%"],
    ["最大单笔收益", "+375.6%"],
]
add_table_slide(prs, "三、A股关键点选股 - 回测结果(2025.06-2026.01)", table_data3)

# 12. A股收益分布
add_content_slide(prs, "三、A股关键点选股 - 收益分布", [
    "【21日收益分布】",
    ">100%: 77个 (0.3%)  █",
    "50%~100%: 517个 (1.8%)  ██",
    "20%~50%: 3,234个 (11.4%)  ████████",
    "10%~20%: 5,158个 (18.2%)  ████████████",
    "5%~10%: 6,247个 (22.0%)  ██████████████",
    "0~5%: 9,706个 (34.3%)  ████████████████████",
    "<0%: 3,397个 (12.0%)  ████████",
    "",
    "【最佳买点案例】",
    "• 嘉美包装(002969): +375.6% (2025-12-03)",
    "• 胜通能源(001331): +322.2% (2025-12-04)",
    "• 合富中国(603122): +256.2% (2025-10-28)"
])

# 13. 实践生产
add_content_slide(prs, "五、实践生产 - 实时监控", [
    "【期货实时信号检测】",
    "• 实时获取1分钟K线数据",
    "• 自动重采样(1min/5min/15min)",
    "• 检测放量K + 2K确认 + 突破买入",
    "• 过滤重复信号",
    "",
    "【A股每日选股】",
    "• 每日收盘后同步数据",
    "• 筛选次日潜在买点",
    "• 输出近期买点(30天内)",
    "• 推送今日信号股票"
])

# 14. 输出文件
add_content_slide(prs, "五、实践生产 - 输出文件", [
    "【期货输出】",
    "• futures_trades_*.csv - 逐笔交易记录",
    "• futures_backtest_summary.csv - 回测统计",
    "",
    "【A股输出】",
    "• stock_recent_buy_points.csv - 近期买点(30天) 2,375个",
    "• stock_today_signals.csv - 今日信号 40只",
    "• stock_all_buy_points_simple.csv - 所有买点 31,568个",
    "• stock_buy_points_detailed.csv - 详细数据+卖点分析",
    "• stock_buy_points_summary.csv - 按股票汇总统计"
])

# 15. 总结
add_content_slide(prs, "六、总结与展望", [
    "【系统能力】",
    "✅ 期货/A股数据自动同步",
    "✅ 主力合约智能筛选",
    "✅ 多周期K线支持",
    "✅ 放量K线检测 + 趋势确认",
    "✅ 止盈止损风控",
    "✅ 历史回测验证",
    "✅ 实时信号推送",
    "",
    "【下一步优化】",
    "→ 期货: 增加30min/60min周期",
    "→ A股: 增加基本面过滤",
    "→ 系统: 接入实盘交易API"
])

# 16. 结束页
add_title_slide(prs, "谢谢!", "量化交易系统 v1.0 - 2026-03-05")

# 保存PPT
output_path = "/Users/qiyue/Desktop/test/claw/quant/docs/量化交易系统汇报.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
